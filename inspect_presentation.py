from pptx import Presentation
import os
p = Presentation('Dashboard/Business_Insights_Presentation.pptx')
print('slides_count=', len(p.slides))
for i, slide in enumerate(p.slides, start=1):
    title = '(no title)'
    try:
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text
    except Exception:
        pass
    pics = 0
    pic_info = []
    for shape in slide.shapes:
        try:
            if hasattr(shape, 'image') and shape.image is not None:
                pics += 1
                pic_info.append(f'image_size={len(shape.image.blob)}')
        except Exception:
            pass
    print(f'slide {i}: title="{title}", pictures={pics}, {pic_info}')

# Also list image files in Images/
img_dir = 'Images'
if os.path.exists(img_dir):
    files = [f for f in os.listdir(img_dir) if f.lower().endswith('.png')]
    print('image_files=', files)
else:
    print('Images folder missing')

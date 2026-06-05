import os
from pptx import Presentation
from pptx.util import Inches

OUTPUT_DIR = "Dashboard"
IMAGES_DIR = "Images"
REPORT_PATH = os.path.join(OUTPUT_DIR, 'Business_Insights_Report.md')
PPTX_PATH = os.path.join(OUTPUT_DIR, 'Business_Insights_Presentation.pptx')
IMAGE_FILES = [
    'monthly_sales_profit.png',
    'region_sales_profit.png',
    'category_profit_margin.png',
    'ship_mode_profit.png',
    'segment_profit_margin.png',
    'top_products_profit.png'
]


def read_report_summary():
    if not os.path.exists(REPORT_PATH):
        return [
            'Business insights report not found.',
            'Run sales_dashboard_analysis.py first to generate the report.'
        ]

    with open(REPORT_PATH, 'r', encoding='utf-8') as f:
        lines = [line.strip() for line in f if line.strip()]

    # Keep relevant summary lines only
    summary_lines = []
    for line in lines[:30]:
        if line.startswith('#'):
            continue
        summary_lines.append(line)
        if len(summary_lines) >= 12:
            break
    return summary_lines


def add_title_slide(prs):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = 'Sales Analytics Dashboard'
    if slide.placeholders:
        subtitle = slide.placeholders[1]
        subtitle.text = 'Superstore Sales Analysis — Client-ready Insights'


def add_summary_slide(prs, lines):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = 'Executive Summary'
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for line in lines:
        p = body.add_paragraph()
        p.text = line
        p.level = 0


def add_image_slide(prs, title_text, image_filename):
    img_path = os.path.join(IMAGES_DIR, image_filename)
    if not os.path.exists(img_path):
        print(f"Skipping missing image: {img_path}")
        return
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    placeholder = slide.shapes.placeholders[1]
    placeholder.text = ''
    left = Inches(0.5)
    top = Inches(1.3)
    height = Inches(5.5)
    slide.shapes.add_picture(img_path, left, top, height=height)


def add_next_steps_slide(prs):
    lines = [
        '1. Open the presentation and review visuals.',
        '2. Open Power BI Desktop and import superstore_clean.csv.',
        '3. Create interactive visuals using the generated charts and insights.',
        '4. Use the report findings to explain recommendations to a business owner.',
        '5. Update the GitHub repo with screenshots and project notes.'
    ]
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = 'Next Steps'
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for line in lines:
        p = body.add_paragraph()
        p.text = line
        p.level = 0


def build_presentation():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    prs = Presentation()
    add_title_slide(prs)
    summary_lines = read_report_summary()
    add_summary_slide(prs, summary_lines)
    for image in IMAGE_FILES:
        title_text = image.replace('_', ' ').replace('.png', '').title()
        add_image_slide(prs, title_text, image)
    add_next_steps_slide(prs)
    prs.save(PPTX_PATH)
    print(f'Saved presentation: {PPTX_PATH}')


if __name__ == '__main__':
    build_presentation()
